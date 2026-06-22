#!/usr/bin/env python3
"""精简提取: 把 Office 文件每行只保留非空单元格, 去掉大量空格/空管道噪音, 便于读材料牌号。
用法: python3 dump_clean.py "<文件或文件夹>"
支持 .xlsx .xlsm .xls .docx .pptx ; PDF 跳过(用 Read 工具)。"""
import os, sys, warnings
warnings.filterwarnings("ignore")

def cell(v):
    if v is None: return ""
    s = str(v).strip()
    return s

def dump_xlsx(p):
    import openpyxl
    wb = openpyxl.load_workbook(p, data_only=True)
    out = []
    for ws in wb.worksheets:
        rows = []
        for r in ws.iter_rows(values_only=True):
            vals = [cell(c) for c in r]
            vals = [v for v in vals if v]
            if vals:
                rows.append(" | ".join(vals))
        if rows:
            out.append(f"### sheet: {ws.title} ({len(rows)} rows) ###")
            out.extend(rows)
    return "\n".join(out)

def dump_xls(p):
    import xlrd
    wb = xlrd.open_workbook(p)
    out = []
    for ws in wb.sheets():
        rows = []
        for r in range(ws.nrows):
            vals = [cell(ws.cell_value(r, c)) for c in range(ws.ncols)]
            vals = [v for v in vals if v]
            if vals:
                rows.append(" | ".join(vals))
        if rows:
            out.append(f"### sheet: {ws.name} ({len(rows)} rows) ###")
            out.extend(rows)
    return "\n".join(out)

def dump_docx(p):
    import docx
    d = docx.Document(p)
    out = [x.text.strip() for x in d.paragraphs if x.text.strip()]
    for t in d.tables:
        for row in t.rows:
            vals = [c.text.strip() for c in row.cells if c.text.strip()]
            if vals:
                out.append(" | ".join(vals))
    return "\n".join(out)

def dump_pptx(p):
    from pptx import Presentation
    out = []
    for i, s in enumerate(Presentation(p).slides):
        txts = []
        for sh in s.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                txts.append(sh.text_frame.text.strip())
            if sh.has_table:
                for row in sh.table.rows:
                    vals = [c.text.strip() for c in row.cells if c.text.strip()]
                    if vals: txts.append(" | ".join(vals))
        if txts:
            out.append(f"--- slide {i+1} ---\n" + "\n".join(txts))
    return "\n".join(out)

H = {"xlsx": dump_xlsx, "xlsm": dump_xlsx, "xls": dump_xls, "docx": dump_docx, "pptx": dump_pptx}

def files(path):
    if os.path.isfile(path):
        yield path
    else:
        for root, _, fs in os.walk(path):
            for f in sorted(fs):
                if not f.startswith("."):
                    yield os.path.join(root, f)

for p in files(sys.argv[1]):
    ext = p.lower().rsplit(".", 1)[-1] if "." in p else ""
    print("\n" + "="*90 + f"\nFILE: {p}\n" + "="*90)
    if ext == "pdf":
        print("[PDF - 用 Read 工具]")
    elif ext in H:
        try: print(H[ext](p))
        except Exception as e: print(f"ERROR: {e}")
    else:
        print(f"[skip .{ext}]")
