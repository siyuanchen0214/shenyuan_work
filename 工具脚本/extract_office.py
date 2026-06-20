#!/usr/bin/env python3
"""
通用询价资料文本提取器 —— 把一个文件夹里所有 Office 文件 dump 成纯文本，便于报价时一次读全。

用法:
    python3 extract_office.py "<文件夹或单个文件路径>" [--max-rows 120]

支持: .pptx .docx .xls .xlsx .xlsm
说明:
  - PDF 不在此处理（用 Claude 的 Read 工具直接读 PDF，可看图纸/3D 缩略图，效果更好）。
  - xlsx/xlsm 用 data_only（取公式计算值），并跳过 openpyxl read_only 模式（它对部分
    带样式的表会报 'ReadOnlyWorksheet has no attribute dimensions'）。
  - 自动跳过纯空 sheet 和已知的二进制垃圾 sheet（如 PIS 模板里的 FontSheet）。
依赖: python-pptx, python-docx, openpyxl, xlrd  (均为常用库，环境已装)
"""
import os, sys, glob, warnings
warnings.filterwarnings("ignore")

SKIP_SHEETS = {"FontSheet"}          # 已知二进制/无意义 sheet
MAX_ROWS_DEFAULT = 120


def _collapse(line: str) -> str:
    return " ".join(line.split())


def dump_pptx(p):
    from pptx import Presentation
    out = []
    for i, s in enumerate(Presentation(p).slides):
        txts = []
        for sh in s.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                txts.append(sh.text_frame.text.strip())
            if sh.has_table:
                for row in sh.table.rows:
                    txts.append(" | ".join(c.text.strip() for c in row.cells))
        if txts:
            out.append(f"--- slide {i+1} ---\n" + "\n".join(txts))
    return "\n".join(out)


def dump_docx(p):
    import docx
    d = docx.Document(p)
    out = [para.text for para in d.paragraphs if para.text.strip()]
    for t in d.tables:
        for row in t.rows:
            out.append(" | ".join(c.text.strip() for c in row.cells))
    return "\n".join(out)


def dump_xlsx(p, max_rows):
    import openpyxl
    wb = openpyxl.load_workbook(p, data_only=True)   # 注意: 不用 read_only
    out = []
    for ws in wb.worksheets:
        if ws.title in SKIP_SHEETS:
            continue
        rows = [r for r in ws.iter_rows(values_only=True)
                if any(c is not None and str(c).strip() for c in r)]
        if not rows:
            continue
        out.append(f"### sheet: {ws.title} ({len(rows)} nonempty rows) ###")
        for r in rows[:max_rows]:
            line = _collapse(" | ".join("" if c is None else str(c).strip() for c in r))
            if line.strip(" |"):
                out.append(line)
        if len(rows) > max_rows:
            out.append("...(truncated)")
    return "\n".join(out)


def dump_xls(p, max_rows):
    import xlrd
    wb = xlrd.open_workbook(p)
    out = []
    for ws in wb.sheets():
        if ws.nrows == 0:
            continue
        out.append(f"### sheet: {ws.name} ({ws.nrows}x{ws.ncols}) ###")
        for r in range(min(ws.nrows, max_rows)):
            vals = [ws.cell_value(r, c) for c in range(ws.ncols)]
            if any(str(v).strip() for v in vals):
                out.append(_collapse(" | ".join(str(v) for v in vals)))
    return "\n".join(out)


HANDLERS = {
    "pptx": lambda p, m: dump_pptx(p),
    "docx": lambda p, m: dump_docx(p),
    "xlsx": dump_xlsx, "xlsm": dump_xlsx,
    "xls":  dump_xls,
}


def iter_files(path):
    if os.path.isfile(path):
        yield path
    else:
        for root, _, files in os.walk(path):
            for f in sorted(files):
                if not f.startswith("."):
                    yield os.path.join(root, f)


def main():
    if len(sys.argv) < 2:
        print(__doc__); sys.exit(1)
    path = sys.argv[1]
    max_rows = MAX_ROWS_DEFAULT
    if "--max-rows" in sys.argv:
        max_rows = int(sys.argv[sys.argv.index("--max-rows") + 1])
    for p in iter_files(path):
        ext = p.lower().rsplit(".", 1)[-1] if "." in p else ""
        print("\n" + "=" * 90 + f"\nFILE: {p}\n" + "=" * 90)
        if ext == "pdf":
            print("[PDF —— 用 Claude Read 工具直接读，可看图纸缩略图]")
        elif ext in HANDLERS:
            try:
                print(HANDLERS[ext](p, max_rows))
            except Exception as e:
                print(f"ERROR reading: {e}")
        else:
            print(f"[skip: .{ext}]")


if __name__ == "__main__":
    main()
