#!/usr/bin/env python3
"""双语标注工具库:在原文(英文)后/下加中文,原文保留。供各 build_*.py 调用。
- docx: 在每个英文段落后插入一段中文(灰色);表格单元格内英文下加中文。
- xlsx/xlsm: 单元格内「英文\n中文」,开启自动换行。跳过 FontSheet 等垃圾sheld。
- pptx: 文本框内英文下加一行中文;纯图片页加中文文本框。
翻译靠传入的 dict {英文原串: 中文}。dict 里没有的串(如牌号/代码)保持原样不动。
材料牌号/件号/代码 不放进 dict => 不翻译。"""
import copy, re
from docx.oxml.ns import qn
from docx.shared import Pt, RGBColor
import docx


# ---------------- docx ----------------
def _insert_para_after(paragraph, text):
    new_p = copy.deepcopy(paragraph._p)
    # 清空 runs
    for r in list(new_p.findall(qn('w:r'))):
        new_p.remove(r)
    paragraph._p.addnext(new_p)
    np = docx.text.paragraph.Paragraph(new_p, paragraph._parent)
    run = np.add_run("【中】" + text)
    run.italic = True
    run.font.color.rgb = RGBColor(0x1F, 0x6F, 0xC0)
    run.font.size = Pt(10)
    return np


def annotate_docx(path, out, tmap):
    d = docx.Document(path)
    # 段落:倒序处理避免插入影响索引
    for p in list(d.paragraphs):
        key = p.text.strip()
        if key and key in tmap:
            _insert_para_after(p, tmap[key])
    # 表格单元格
    for t in d.tables:
        for row in t.rows:
            for cell in row.cells:
                key = cell.text.strip()
                if key and key in tmap:
                    run = cell.paragraphs[-1].add_run("\n【中】" + tmap[key])
                    run.italic = True
                    run.font.color.rgb = RGBColor(0x1F, 0x6F, 0xC0)
    d.save(out)


# ---------------- xlsx / xlsm ----------------
import openpyxl
from openpyxl.styles import Alignment, Font

SKIP_SHEETS = {"FontSheet"}


def annotate_xlsx(path, out, tmap, keep_vba=False):
    wb = openpyxl.load_workbook(path, keep_vba=keep_vba)
    for ws in wb.worksheets:
        if ws.title in SKIP_SHEETS:
            continue
        for row in ws.iter_rows():
            for cell in row:
                v = cell.value
                if not isinstance(v, str):
                    continue
                key = v.strip()
                if key in tmap and tmap[key]:
                    cell.value = v + "\n【中】" + tmap[key]
                    al = cell.alignment
                    cell.alignment = Alignment(wrap_text=True,
                                               horizontal=al.horizontal,
                                               vertical=al.vertical or "top")
    wb.save(out)


# ---------------- xls (旧版,转 xlsx 双语) ----------------
import xlrd


def annotate_xls_to_xlsx(path, out, tmap):
    rb = xlrd.open_workbook(path, formatting_info=False)
    wb = openpyxl.Workbook()
    wb.remove(wb.active)
    for sh in rb.sheets():
        ws = wb.create_sheet(title=sh.name[:31] or "Sheet")
        for r in range(sh.nrows):
            for c in range(sh.ncols):
                v = sh.cell_value(r, c)
                if v == "":
                    continue
                if isinstance(v, str) and v.strip() in tmap and tmap[v.strip()]:
                    out_v = v + "\n【中】" + tmap[v.strip()]
                    ws.cell(r + 1, c + 1, out_v).alignment = Alignment(wrap_text=True, vertical="top")
                else:
                    ws.cell(r + 1, c + 1, v)
        ws.column_dimensions  # noop
    wb.save(out)


# ---------------- pptx ----------------
from pptx import Presentation
from pptx.util import Emu, Pt as PPt
from pptx.dml.color import RGBColor as PRGB


def annotate_pptx(path, out, tmap, slide_overlays=None):
    """tmap: 文本框文字→中文,逐段在英文下加中文。
    slide_overlays: {slide_index(0基): '整段中文'} 用于纯图片页叠加中文文本框。"""
    prs = Presentation(path)

    def do_tf(tf):
        for para in tf.paragraphs:
            txt = "".join(r.text for r in para.runs).strip()
            if txt and txt in tmap and tmap[txt]:
                run = para.add_run()
                run.text = "  【中】" + tmap[txt]
                run.font.size = PPt(10)
                run.font.color.rgb = PRGB(0x1F, 0x6F, 0xC0)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                do_tf(shape.text_frame)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        do_tf(cell.text_frame)
    if slide_overlays:
        slides = list(prs.slides)
        for idx, text in slide_overlays.items():
            if idx < len(slides):
                s = slides[idx]
                tb = s.shapes.add_textbox(Emu(228600), Emu(228600),
                                          Emu(8000000), Emu(1200000))
                tf = tb.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                r = p.add_run(); r.text = "【中文翻译】" + text
                r.font.size = PPt(12); r.font.color.rgb = PRGB(0xC0, 0x1F, 0x1F)
    prs.save(out)
